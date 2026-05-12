"""
Exit-Game PowerPoint Generator
Erstellt eine professionelle 16:9 Präsentation mit dunklem Theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── Farb-Palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # Haupt-Hintergrund
BG_CARD      = RGBColor(0x16, 0x1B, 0x22)   # Karten-Hintergrund
BG_CARD2     = RGBColor(0x1C, 0x23, 0x2E)   # Heller Karten-Hintergrund
ACCENT_GREEN = RGBColor(0x39, 0xD3, 0x53)   # Primär-Akzent Grün
ACCENT_BLUE  = RGBColor(0x58, 0xA6, 0xFF)   # Sekundär-Akzent Blau
ACCENT_GOLD  = RGBColor(0xF0, 0xC0, 0x40)   # Warn / HIGH-Priority
ACCENT_RED   = RGBColor(0xFF, 0x5C, 0x5C)   # Fehler / Rot
ACCENT_TEAL  = RGBColor(0x00, 0xBF, 0xC4)   # Arduino Teal
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xC9, 0xD1, 0xD9)
TEXT_GREY    = RGBColor(0x8B, 0x94, 0x9E)
TEXT_DARK    = RGBColor(0x0D, 0x11, 0x17)
LINE_COLOR   = RGBColor(0x30, 0x36, 0x3D)

# ── Hilfsfunktionen ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_text(slide, text, x, y, w, h, font_size=16, bold=False, color=TEXT_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_label(slide, text, x, y, w=0.18, h=0.28, bg=ACCENT_GREEN, fg=TEXT_DARK):
    r = add_rect(slide, x, y, w, h, bg)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.color.rgb = fg
    return r

def slide_header(slide, title, subtitle=None, accent=ACCENT_GREEN):
    # Hintergrundfarbe
    set_bg(slide, BG_DARK)
    # Linker Akzentbalken
    add_rect(slide, 0, 0, 0.06, 7.5, accent)
    # Titel
    add_text(slide, title, 0.35, 0.22, 12.5, 0.65,
             font_size=28, bold=True, color=TEXT_WHITE)
    # Trennlinie
    add_rect(slide, 0.35, 0.92, 12.5, 0.025, accent)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.96, 12.5, 0.4,
                 font_size=13, bold=False, color=TEXT_GREY)

def add_card(slide, x, y, w, h, color=BG_CARD):
    r = add_rect(slide, x, y, w, h, color)
    return r

def add_bullet_card(slide, x, y, w, h, title, bullets, accent=ACCENT_GREEN,
                    title_size=13, bullet_size=11, bg=BG_CARD):
    add_card(slide, x, y, w, h, bg)
    add_rect(slide, x, y, 0.045, h, accent)
    add_text(slide, title, x + 0.12, y + 0.08, w - 0.2, 0.32,
             font_size=title_size, bold=True, color=TEXT_WHITE)
    for i, b in enumerate(bullets):
        add_text(slide, f"• {b}", x + 0.12, y + 0.42 + i * 0.32, w - 0.2, 0.3,
                 font_size=bullet_size, color=TEXT_LIGHT)

def add_icon_row(slide, items, y, col_w=2.5, start_x=0.5, accent=ACCENT_GREEN):
    for i, (icon, label, sub) in enumerate(items):
        x = start_x + i * col_w
        add_rect(slide, x, y, col_w - 0.15, 1.5, BG_CARD)
        add_rect(slide, x, y, col_w - 0.15, 0.07, accent)
        add_text(slide, icon, x + 0.1, y + 0.12, col_w - 0.35, 0.38,
                 font_size=20, bold=True, color=accent)
        add_text(slide, label, x + 0.1, y + 0.55, col_w - 0.35, 0.35,
                 font_size=11, bold=True, color=TEXT_WHITE)
        if sub:
            add_text(slide, sub, x + 0.1, y + 0.92, col_w - 0.35, 0.45,
                     font_size=9, color=TEXT_GREY)

# ── Präsentation erstellen ────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # Blank

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 1 – Titel
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG_DARK)

# Großer Hintergrundakzent links
add_rect(s, 0, 0, 4.2, 7.5, BG_CARD)
add_rect(s, 0, 0, 0.09, 7.5, ACCENT_GREEN)

# Grid-Dekor rechts (Punkte-Matrix simuliert)
for row in range(6):
    for col in range(10):
        add_rect(s, 4.6 + col * 0.88, 0.5 + row * 1.05, 0.06, 0.06, LINE_COLOR)

# Haupt-Titel
add_text(s, "EXIT-GAME", 0.35, 1.4, 4.0, 0.95,
         font_size=42, bold=True, color=ACCENT_GREEN)
add_text(s, "Projektumsetzung", 0.35, 2.35, 4.0, 0.7,
         font_size=28, bold=True, color=TEXT_WHITE)

# Trennlinie
add_rect(s, 0.35, 3.12, 3.6, 0.04, ACCENT_GREEN)

# Untertitel
add_text(s,
         "Unity-basiertes Escape-Game\nmit Arduino-Interaktion,\nNPC-Dialogen und\ntestgetriebener Entwicklung",
         0.35, 3.25, 3.7, 1.6,
         font_size=13, color=TEXT_LIGHT)

# Tech-Tags
for i, tag in enumerate(["Unity 2022+", "C# / .NET", "Arduino", "TDD"]):
    add_rect(s, 0.35 + i * 0.95, 5.6, 0.85, 0.3, ACCENT_BLUE)
    add_text(s, tag, 0.37 + i * 0.95, 5.62, 0.83, 0.26,
             font_size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# Rechts: Karten mit Level-Nummern
level_colors = [ACCENT_GREEN, ACCENT_BLUE, ACCENT_TEAL, ACCENT_GOLD,
                RGBColor(0xAA, 0x66, 0xFF), ACCENT_RED]
level_names  = ["Zelle", "Abstellr.", "Bibliothek", "Computer", "Werkstatt", "Tor"]
for i in range(6):
    col = i % 3
    row = i // 3
    cx = 5.1 + col * 2.7
    cy = 1.8 + row * 2.4
    add_rect(s, cx, cy, 2.4, 2.0, BG_CARD2)
    add_rect(s, cx, cy, 2.4, 0.06, level_colors[i])
    add_text(s, f"L{i+1}", cx + 0.15, cy + 0.12, 0.7, 0.65,
             font_size=32, bold=True, color=level_colors[i])
    add_text(s, f"Level {i+1}", cx + 0.15, cy + 0.78, 2.0, 0.32,
             font_size=12, bold=True, color=TEXT_WHITE)
    add_text(s, level_names[i], cx + 0.15, cy + 1.12, 2.0, 0.28,
             font_size=10, color=TEXT_GREY)

add_text(s, "TH Köln · Sommersemester 2025", 0.35, 7.05, 5.0, 0.35,
         font_size=9, color=TEXT_GREY)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 2 – Projektüberblick
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Projektüberblick", "Was bauen wir – und warum?")

# 3 Spalten-Karten
cards = [
    ("🎮", "6 Level", "Escape-Game",
     ["Prison Cell → Storage Room\n→ Library → Computer\n→ Workshop → Final Gate",
      "Aufeinander aufbauende\nPuzzle-Mechaniken",
      "Jedes Level ein eigener\nScene-Build"]),
    ("🧩", "Hybrid-Ansatz", "3D + UI + Hardware",
     ["3D-Szenen (Level 1–2, 5–6)\nfür immersives Erlebnis",
      "UI-Panels (Level 3–6)\nfür Puzzle-Fokus",
      "Arduino-Sensoren für echte\nphysische Interaktion"]),
    ("⚙️", "Technisch", "Modular & Testbar",
     ["Singleton-Architektur\n(GameManager, ArduinoBridge)",
      "Event-driven Kommunikation\nzwischen Komponenten",
      "Test-driven Development\n(EditMode + PlayMode Tests)"]),
]
for i, (icon, title, sub, bullets) in enumerate(cards):
    x = 0.35 + i * 4.3
    add_rect(s, x, 1.55, 4.1, 5.55, BG_CARD)
    add_rect(s, x, 1.55, 4.1, 0.06, ACCENT_GREEN if i == 0 else ACCENT_BLUE if i == 1 else ACCENT_TEAL)
    add_text(s, icon, x + 0.18, 1.65, 0.7, 0.55, font_size=24, color=TEXT_WHITE)
    add_text(s, title, x + 0.18, 2.22, 3.6, 0.38,
             font_size=16, bold=True, color=TEXT_WHITE)
    add_text(s, sub, x + 0.18, 2.62, 3.6, 0.3,
             font_size=10, color=TEXT_GREY)
    add_rect(s, x + 0.18, 3.0, 3.5, 0.03, LINE_COLOR)
    for j, b in enumerate(bullets):
        add_text(s, f"▸  {b}", x + 0.18, 3.12 + j * 0.95, 3.6, 0.85,
                 font_size=10, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 3 – Technologiestack
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Technologiestack", "Eingesetzte Technologien und Frameworks")

tech_items = [
    ("🎯", "Unity 2022+", "Game Engine",
     "3D-Szenen, UI-Panels, Physics, Input System, Play/Edit Mode Tests"),
    ("💎", "C# / .NET", "Sprache",
     "MonoBehaviour Lifecycle, Coroutines, Async Threads, LINQ"),
    ("🖥️", "TextMeshPro + Unity UI", "Rendering",
     "Canvas-basierte Panels, Typewriter-Dialoge, Responsive Buttons"),
    ("🎮", "Unity Input System", "Eingabe",
     "Keyboard, Mouse, Gamepad – kompatibel mit allen Level-Steuerungen"),
    ("🔌", "Arduino + SerialPort", "Hardware",
     "System.IO.Ports, dedizierter Background Thread, Hex-Protokoll"),
    ("🧪", "NUnit / Unity Test Runner", "Testing",
     "EditMode & PlayMode Tests, Mindestabdeckung pro Level definiert"),
]

for i, (icon, name, cat, desc) in enumerate(tech_items):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.3
    y = 1.55 + row * 2.75
    add_rect(s, x, y, 4.1, 2.5, BG_CARD)
    # Farbiger linker Rand
    col_accent = [ACCENT_GREEN, ACCENT_BLUE, ACCENT_TEAL,
                  ACCENT_GOLD, ACCENT_RED, RGBColor(0xAA, 0x66, 0xFF)][i]
    add_rect(s, x, y, 0.045, 2.5, col_accent)
    add_text(s, icon, x + 0.15, y + 0.12, 0.55, 0.5, font_size=22, color=TEXT_WHITE)
    add_text(s, name, x + 0.75, y + 0.12, 3.2, 0.38,
             font_size=14, bold=True, color=TEXT_WHITE)
    add_text(s, cat, x + 0.75, y + 0.5, 3.2, 0.28,
             font_size=9, bold=True, color=col_accent)
    add_text(s, desc, x + 0.15, y + 0.95, 3.8, 1.3,
             font_size=9.5, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 4 – Rollenmodell / Agententeam
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Virtuelles Agenten-Team", "Interne Rollenverteilung für koordinierte Entwicklung")

roles = [
    ("🏛️", "Architect", ACCENT_BLUE,
     "Scene hierarchy\nData flow\nSystem design decisions"),
    ("⚡", "Unity-Expert", ACCENT_GREEN,
     "MonoBehaviour lifecycle\nPhysics & UI\nAnimation & Input System"),
    ("📡", "Arduino-Bridge", ACCENT_TEAL,
     "Serial communication layer\nSensor protocol\nAsync data handling"),
    ("🔬", "QA-Engineer", ACCENT_GOLD,
     "Edit/Play Mode Tests\nRegression checks\nPerformance profiling"),
    ("🎯", "Head of Dev", ACCENT_RED,
     "Koordiniert alle Rollen\nOwns Git commits\nFinal review"),
]

# Zentrum: Head of Dev
add_rect(s, 4.9, 2.8, 3.5, 2.4, BG_CARD2)
add_rect(s, 4.9, 2.8, 3.5, 0.07, ACCENT_RED)
add_text(s, "🎯", 5.6, 2.92, 0.6, 0.55, font_size=24, color=TEXT_WHITE)
add_text(s, "HEAD OF DEVELOPMENT", 6.25, 2.95, 2.0, 0.38,
         font_size=10, bold=True, color=ACCENT_RED)
add_text(s, "Koordiniert alle Rollen\nOwns Git commits & Final Review", 5.05, 3.42, 3.2, 0.75,
         font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# 4 äußere Rollen
outer_positions = [(0.25, 1.4), (9.55, 1.4), (0.25, 4.65), (9.55, 4.65)]
outer_roles = [roles[0], roles[1], roles[2], roles[3]]

for (rx, ry), (icon, name, accent, desc) in zip(outer_positions, outer_roles):
    add_rect(s, rx, ry, 3.1, 2.2, BG_CARD)
    add_rect(s, rx, ry, 3.1, 0.055, accent)
    add_text(s, icon, rx + 0.15, ry + 0.12, 0.55, 0.48, font_size=20, color=TEXT_WHITE)
    add_text(s, name.upper(), rx + 0.75, ry + 0.12, 2.2, 0.35,
             font_size=12, bold=True, color=accent)
    add_text(s, desc, rx + 0.15, ry + 0.62, 2.8, 1.45,
             font_size=9.5, color=TEXT_LIGHT)

# Verbindungslinien (als dünne Rechtecke)
# Mitte zu links oben / rechts oben / links unten / rechts unten
connector_coords = [
    (3.35, 2.55, 1.6, 0.04),  # → links oben
    (8.75, 2.55, 1.6, 0.04),  # → rechts oben (von Mitte nach rechts)
    (3.35, 5.08, 1.6, 0.04),  # → links unten
    (8.75, 5.08, 1.6, 0.04),  # → rechts unten
]
for (cx, cy, cw, ch) in connector_coords:
    add_rect(s, cx, cy, cw, ch, LINE_COLOR)

add_text(s, "Jede Rolle bündelt Expertise → alle arbeiten am selben Codebase",
         0.35, 7.1, 12.5, 0.3, font_size=9, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 5 – Projektarchitektur
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Projektarchitektur", "Szenenstruktur, Singletons und Kommunikationsfluss")

# Architektur-Diagramm als Boxen
# Zeile 1: Szenen
add_rect(s, 0.35, 1.6, 12.7, 0.35, RGBColor(0x1A, 0x22, 0x30))
add_text(s, "SZENEN", 0.55, 1.65, 3.0, 0.25,
         font_size=9, bold=True, color=TEXT_GREY)

scene_boxes = [
    ("Level1.unity", "Prison Cell", ACCENT_GREEN),
    ("Level2.unity", "Storage Room", ACCENT_BLUE),
    ("GameLevel.unity", "Level 3–6 Panels", ACCENT_TEAL),
    ("SampleScene.unity", "Sandbox / Testing", TEXT_GREY),
]
for i, (name, sub, col) in enumerate(scene_boxes):
    x = 0.35 + i * 3.18
    add_rect(s, x, 2.0, 3.0, 0.95, BG_CARD)
    add_rect(s, x, 2.0, 3.0, 0.055, col)
    add_text(s, name, x + 0.1, 2.08, 2.8, 0.32,
             font_size=10, bold=True, color=TEXT_WHITE)
    add_text(s, sub, x + 0.1, 2.42, 2.8, 0.42,
             font_size=9, color=col)

# Pfeil nach unten (Singletons)
add_rect(s, 6.3, 3.02, 0.04, 0.45, LINE_COLOR)
add_text(s, "▼", 6.15, 3.4, 0.4, 0.3, font_size=11, color=TEXT_GREY)

# Zeile 2: Singletons
add_rect(s, 0.35, 3.5, 12.7, 0.35, RGBColor(0x1A, 0x22, 0x30))
add_text(s, "CORE SINGLETONS (DontDestroyOnLoad)", 0.55, 3.55, 6.0, 0.25,
         font_size=9, bold=True, color=TEXT_GREY)

singletons = [
    ("GameManager", "Level-Progression\nCompleteCurrentLevel()\nPanel-Activation", ACCENT_GREEN, 0.35),
    ("BigYahuDialogSystem", "Typewriter-Dialog\nPortrait & Audio\nCallback-Chain", ACCENT_BLUE, 4.45),
    ("ArduinoBridge", "Serial / TCP\nBackground Thread\nEvent-Dispatch", ACCENT_TEAL, 8.55),
]
for name, desc, col, x in singletons:
    add_rect(s, x, 3.9, 3.85, 1.5, BG_CARD2)
    add_rect(s, x, 3.9, 3.85, 0.055, col)
    add_text(s, name, x + 0.12, 3.98, 3.5, 0.35,
             font_size=12, bold=True, color=col)
    add_text(s, desc, x + 0.12, 4.38, 3.5, 1.0,
             font_size=9.5, color=TEXT_LIGHT)

# Pfeil nach unten (Level-Scripts)
add_rect(s, 6.3, 5.45, 0.04, 0.3, LINE_COLOR)
add_text(s, "▼", 6.15, 5.65, 0.4, 0.28, font_size=11, color=TEXT_GREY)

# Zeile 3: Level-Scripts + UI + Tests
add_rect(s, 0.35, 5.78, 12.7, 0.32, RGBColor(0x1A, 0x22, 0x30))
add_text(s, "LEVEL-SCRIPTS + UI PANELS + TESTS", 0.55, 5.82, 8.0, 0.22,
         font_size=9, bold=True, color=TEXT_GREY)

bottom_boxes = [
    ("Level Scripts\nLevel{N}_{Desc}.cs", ACCENT_GREEN, 0.35, 3.85),
    ("UI Panels\nCanvas / TMP / Button", ACCENT_BLUE, 4.45, 3.85),
    ("Tests\nEditMode + PlayMode", ACCENT_GOLD, 8.55, 3.85),
]
for name, col, x, w in bottom_boxes:
    add_rect(s, x, 6.15, w, 1.1, BG_CARD)
    add_rect(s, x, 6.15, w, 0.045, col)
    add_text(s, name, x + 0.12, 6.22, w - 0.2, 0.8,
             font_size=10, color=TEXT_WHITE)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 6 – Level-Umsetzung (Tabelle)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Level-Umsetzung", "6 Level — je eigenständig, vollständig konfiguriert")

# Tabelle manuell mit Rechtecken und Text
headers = ["Level", "Szene / Panel", "Script", "Arduino Hook"]
col_xs  = [0.35, 1.85, 4.35, 8.85]
col_ws  = [1.45, 2.45, 4.45, 4.0]

# Header-Zeile
header_y = 1.55
for j, (hdr, x, w) in enumerate(zip(headers, col_xs, col_ws)):
    add_rect(s, x, header_y, w - 0.05, 0.42, BG_CARD2)
    add_text(s, hdr, x + 0.08, header_y + 0.07, w - 0.2, 0.3,
             font_size=10, bold=True, color=ACCENT_GREEN)

# Daten-Zeilen
level_colors_row = [ACCENT_GREEN, ACCENT_BLUE, ACCENT_TEAL,
                    ACCENT_GOLD, RGBColor(0xAA, 0x66, 0xFF), ACCENT_RED]
rows = [
    ["L1 – Zelle",       "Level1.unity",     "Level1_Cell.cs\nToiletInteraction.cs\nNumpadController.cs", "—"],
    ["L2 – Abstellr.",   "Level2.unity",     "Level2_DustWall.cs",       "Humidity Sensor\n(Blow Detection)"],
    ["L3 – Bibliothek",  "GameLevel (Panel)","Level3_ColorCode.cs",      "Color Sensor\n(Sequence Scan)"],
    ["L4 – Computer",    "GameLevel (Panel)","Level4_StealthMinigame.cs","—"],
    ["L5 – Werkstatt",   "GameLevel (Panel)","Level5_Breadboard.cs",     "Breadboard I/O\n(Circuit State)"],
    ["L6 – Tor",         "GameLevel (Panel)","Level6_FinalGate.cs",      "Temperature Sensor\n(Heat Detection)"],
]
row_h = 0.73
for i, row in enumerate(rows):
    ry = header_y + 0.47 + i * row_h
    bg = BG_CARD if i % 2 == 0 else BG_CARD2
    for j, (cell, x, w) in enumerate(zip(row, col_xs, col_ws)):
        add_rect(s, x, ry, w - 0.05, row_h - 0.04, bg)
        if j == 0:
            # Level-Nummer mit Farb-Indikator
            add_rect(s, x, ry, 0.04, row_h - 0.04, level_colors_row[i])
        add_text(s, cell, x + 0.1, ry + 0.08, w - 0.18, row_h - 0.15,
                 font_size=9.5 if j != 2 else 8.5, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 7 – Arduino-Integration
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Arduino-Integration", "Asynchrone Hardware-Kommunikation — sicher & robuste")

# Fluss-Diagramm oben
boxes_flow = [
    ("Arduino\nHardware", ACCENT_TEAL, 0.35),
    ("Background\nThread", ACCENT_BLUE, 3.05),
    ("ConcurrentQueue\n(thread-safe)", ACCENT_GREEN, 5.75),
    ("Unity Main Thread\nUpdate()", ACCENT_GOLD, 8.55),
]
flow_y = 1.55
for i, (label, col, x) in enumerate(boxes_flow):
    add_rect(s, x, flow_y, 2.55, 1.05, BG_CARD)
    add_rect(s, x, flow_y, 2.55, 0.055, col)
    add_text(s, label, x + 0.12, flow_y + 0.12, 2.3, 0.75,
             font_size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "→", x + 2.6, flow_y + 0.35, 0.35, 0.4,
                 font_size=16, bold=True, color=TEXT_GREY)

# Protokoll-Box
add_rect(s, 0.35, 2.85, 6.1, 2.55, BG_CARD)
add_rect(s, 0.35, 2.85, 0.045, 2.55, ACCENT_TEAL)
add_text(s, "Nachrichten-Format", 0.5, 2.92, 5.8, 0.35,
         font_size=12, bold=True, color=ACCENT_TEAL)
proto = [
    ('PC → Arduino:', '0x01 "TRIGGER\\n"', ACCENT_GREEN),
    ('Arduino → PC:', '0x10 "HUMIDITY:72\\n"', ACCENT_BLUE),
    ('Text-Format:', '"KEY:5" | "HUMIDITY:72" | "COLOR:RED" | "TEMP:65"', ACCENT_GOLD),
]
for i, (label, value, col) in enumerate(proto):
    add_text(s, label, 0.5, 3.35 + i * 0.55, 1.6, 0.4,
             font_size=9, bold=True, color=TEXT_GREY)
    add_text(s, value, 2.15, 3.35 + i * 0.55, 4.0, 0.4,
             font_size=9.5, bold=True, color=col)

# Command-IDs Tabelle rechts
add_rect(s, 6.7, 2.85, 6.0, 2.55, BG_CARD)
add_rect(s, 6.7, 2.85, 0.045, 2.55, ACCENT_BLUE)
add_text(s, "Command IDs", 6.85, 2.92, 5.7, 0.35,
         font_size=12, bold=True, color=ACCENT_BLUE)
cmds = [
    ("0x01", "Ping / Keep-Alive"),
    ("0x10", "Humidity Sensor (L2)"),
    ("0x11", "Joystick Direction (L2)"),
    ("0x20", "Color Sensor (L3)"),
    ("0x50", "Breadboard State (L5)"),
    ("0x60", "Temperature Sensor (L6)"),
]
for i, (cid, desc) in enumerate(cmds):
    add_rect(s, 6.85, 3.35 + i * 0.34, 0.58, 0.28, ACCENT_BLUE)
    add_text(s, cid, 6.87, 3.37 + i * 0.34, 0.55, 0.22,
             font_size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, desc, 7.55, 3.37 + i * 0.34, 4.9, 0.22,
             font_size=9, color=TEXT_LIGHT)

# Unten: Prinzipien
principles = [
    "Kein Blocking-Read\nim Main Thread",
    "Timeout + Auto-\nReconnect Logik",
    "Arduino-Fallback:\nKeyboard/Mouse",
    "Hex-Protokoll mit\n1-Byte Command-ID",
]
for i, p in enumerate(principles):
    x = 0.35 + i * 3.25
    add_rect(s, x, 5.65, 3.1, 1.6, BG_CARD2)
    add_rect(s, x, 5.65, 3.1, 0.05, ACCENT_TEAL)
    add_text(s, p, x + 0.15, 5.8, 2.8, 1.35,
             font_size=10, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 8 – Coding Standards
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Coding Standards & Leitlinien", "Einheitliche Qualität, Wartbarkeit und Performance")

# Linke Spalte: Design-Prinzipien
add_rect(s, 0.35, 1.55, 6.0, 5.72, BG_CARD)
add_rect(s, 0.35, 1.55, 0.045, 5.72, ACCENT_GREEN)
add_text(s, "Design-Prinzipien", 0.5, 1.62, 5.7, 0.38,
         font_size=13, bold=True, color=ACCENT_GREEN)

design_rules = [
    ("Single Responsibility", "Ein MonoBehaviour = eine Aufgabe.\nKeine God-Objects."),
    ("Event-Driven", "Action / UnityEvent statt direkter Referenzen.\nKopplung minimieren."),
    ("Kein Polling", "Update() nur wenn nötig.\nEvents, Coroutines und Trigger bevorzugen."),
    ("Caching", "GetComponent<>() in Awake/Start cachen.\nNie in Update() aufrufen."),
    ("Puzzle O(n)", "Validierungsschleifen maximal O(n).\nn ist stets ≤ 10."),
]
for i, (rule, desc) in enumerate(design_rules):
    ry = 2.12 + i * 0.98
    add_rect(s, 0.5, ry, 0.28, 0.28, ACCENT_GREEN)
    add_text(s, "✓", 0.51, ry + 0.01, 0.25, 0.25,
             font_size=9, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, rule, 0.88, ry, 5.2, 0.3,
             font_size=10.5, bold=True, color=TEXT_WHITE)
    add_text(s, desc, 0.88, ry + 0.32, 5.2, 0.55,
             font_size=9, color=TEXT_LIGHT)

# Rechte Spalte: Naming Conventions + DOs/DONTs
add_rect(s, 6.65, 1.55, 6.02, 2.65, BG_CARD)
add_rect(s, 6.65, 1.55, 0.045, 2.65, ACCENT_BLUE)
add_text(s, "Naming Conventions", 6.8, 1.62, 5.7, 0.38,
         font_size=13, bold=True, color=ACCENT_BLUE)
names = [
    ("Scripts:", "PascalCase"),
    ("Level Scripts:", "Level{N}_{Description}.cs"),
    ("Editor Builder:", "Build{LevelName}.cs"),
    ("Arduino:", "Assets/Scripts/Arduino/ArduinoBridge.cs"),
    ("Tests:", "Assets/Tests/EditMode/ & PlayMode/"),
]
for i, (k, v) in enumerate(names):
    add_text(s, k, 6.8, 2.12 + i * 0.38, 1.6, 0.3,
             font_size=9, bold=True, color=TEXT_GREY)
    add_text(s, v, 8.45, 2.12 + i * 0.38, 4.0, 0.3,
             font_size=9, color=TEXT_LIGHT)

# Arduino DO/DON'T
add_rect(s, 6.65, 4.42, 6.02, 2.85, BG_CARD)
add_rect(s, 6.65, 4.42, 0.045, 2.85, ACCENT_TEAL)
add_text(s, "Arduino-Regeln", 6.8, 4.49, 5.7, 0.38,
         font_size=13, bold=True, color=ACCENT_TEAL)
dos   = ["Background Thread für Serial Reads",
         "Timeout + Reconnect Logik",
         "Hex Command-ID Prefix"]
donts = ["SerialPort.Read() im Main Thread",
         "Update() für Hardware warten",
         "Blocking Reads"]
for i, txt in enumerate(dos):
    add_rect(s, 6.8, 5.0 + i * 0.38, 0.28, 0.22, ACCENT_GREEN)
    add_text(s, "DO", 6.82, 5.01 + i * 0.38, 0.26, 0.18,
             font_size=6, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, txt, 7.18, 5.0 + i * 0.38, 2.3, 0.3, font_size=8.5, color=TEXT_LIGHT)
for i, txt in enumerate(donts):
    add_rect(s, 9.55, 5.0 + i * 0.38, 0.38, 0.22, ACCENT_RED)
    add_text(s, "DON'T", 9.57, 5.01 + i * 0.38, 0.35, 0.18,
             font_size=6, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, txt, 10.02, 5.0 + i * 0.38, 2.5, 0.3, font_size=8.5, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 9 – TDD & QA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Test-Driven Development & QA", "Jede Logikkomponente braucht Tests — keine Ausnahmen")

# Test-Typen nebeneinander
test_types = [
    ("📝", "EditMode Tests", ACCENT_BLUE,
     ["Pure Puzzle-Logik", "Code-Parsing", "Zustandsmaschinen",
      "Kein Unity-Lifecycle nötig", "Schnell & deterministisch"]),
    ("▶️", "PlayMode Tests", ACCENT_GREEN,
     ["Coroutine-Flows", "Scene Transitions", "UI-Interaktion",
      "Vollständiger Unity-Lifecycle", "NPC-Dialog-Callbacks"]),
    ("🔌", "Hardware-Fallback", ACCENT_TEAL,
     ["Arduino disconnect simulieren", "Keyboard/Mouse Fallback prüfen",
      "arduinoFallback = true\nals SerializeField", "Graceful Degradation"]),
]
for i, (icon, title, col, bullets) in enumerate(test_types):
    x = 0.35 + i * 4.3
    add_rect(s, x, 1.55, 4.1, 4.0, BG_CARD)
    add_rect(s, x, 1.55, 4.1, 0.055, col)
    add_text(s, icon, x + 0.15, 1.65, 0.55, 0.55, font_size=22, color=TEXT_WHITE)
    add_text(s, title, x + 0.75, 1.68, 3.2, 0.38,
             font_size=14, bold=True, color=col)
    add_rect(s, x + 0.15, 2.18, 3.7, 0.03, LINE_COLOR)
    for j, b in enumerate(bullets):
        add_text(s, f"▸  {b}", x + 0.15, 2.3 + j * 0.62, 3.8, 0.55,
                 font_size=9.5, color=TEXT_LIGHT)

# Mindestabdeckung
add_rect(s, 0.35, 5.7, 12.7, 1.58, BG_CARD2)
add_rect(s, 0.35, 5.7, 12.7, 0.05, ACCENT_GOLD)
add_text(s, "Mindest-Testabdeckung pro Level", 0.55, 5.77, 8.0, 0.38,
         font_size=12, bold=True, color=ACCENT_GOLD)
min_tests = [
    ("✅ Correct Solution", "CompleteCurrentLevel() wird korrekt aufgerufen"),
    ("❌ Wrong Solution", "State wird zurückgesetzt, kein Level-Advance"),
    ("🔌 Arduino Disconnect", "Graceful Degradation, Fallback auf Keyboard/Mouse aktiv"),
]
for i, (label, desc) in enumerate(min_tests):
    x = 0.5 + i * 4.25
    add_text(s, label, x, 6.22, 4.0, 0.3,
             font_size=10, bold=True, color=TEXT_WHITE)
    add_text(s, desc, x, 6.58, 4.0, 0.55,
             font_size=9, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 10 – GDD Gaps
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "GDD → Implementierung: Offene Punkte", "Priorisierte Lücken zwischen Game Design Document und aktuellem Stand")

# Header
gap_headers = ["ID", "GDD-Anforderung", "Aktueller Stand", "Priorität"]
gap_hx      = [0.35, 0.95, 5.55, 10.65]
gap_hw      = [0.55, 4.55, 5.05, 2.25]
header_y_g  = 1.55
for hdr, x, w in zip(gap_headers, gap_hx, gap_hw):
    add_rect(s, x, header_y_g, w - 0.04, 0.4, BG_CARD2)
    add_text(s, hdr, x + 0.06, header_y_g + 0.06, w - 0.12, 0.28,
             font_size=10, bold=True, color=ACCENT_GREEN)

gaps = [
    ("G-01",
     "L1: Hex-Code 66A auf Decke sichtbar → Decode zu 1642",
     "Code korrekt implementiert, aber kein Decken-Objekt in-world",
     "HIGH", ACCENT_RED),
    ("G-02",
     "L2: Joshi NPC Narration, nicht Big Yahu",
     "Level2_DustWall.cs nutzt BigYahuDialogSystem",
     "HIGH", ACCENT_RED),
    ("G-03",
     "L2: Joystick-Kombo ↑↑↓↓ öffnet Tür",
     "Kein Directional-Puzzle; versteckte Nummer wird angezeigt",
     "HIGH", ACCENT_RED),
    ("G-04",
     "L3: Helios NPC + Buch-Auswahl (Bibel wählen)",
     "Kein Buch-Auswahl-Mechanismus; Big Yahu spricht",
     "MEDIUM", ACCENT_GOLD),
    ("G-05",
     "L3: Farbsequenz via Arduino (Sensor an Code halten)",
     "Nur Mausklick-Buttons; kein Arduino-Pfad implementiert",
     "MEDIUM", ACCENT_GOLD),
    ("G-06",
     "L6: Bunsenbrenner / Föhn Metapher",
     "Generisches Heat-Button UI ohne Metapher",
     "LOW", ACCENT_BLUE),
]
row_h_g = 0.78
for i, (gid, req, state, prio, pcol) in enumerate(gaps):
    ry = header_y_g + 0.45 + i * row_h_g
    bg = BG_CARD if i % 2 == 0 else BG_CARD2
    for x, w in zip(gap_hx, gap_hw):
        add_rect(s, x, ry, w - 0.04, row_h_g - 0.04, bg)
    add_text(s, gid, gap_hx[0] + 0.05, ry + 0.08, 0.45, 0.55,
             font_size=10, bold=True, color=pcol, align=PP_ALIGN.CENTER)
    add_text(s, req, gap_hx[1] + 0.06, ry + 0.05, gap_hw[1] - 0.12, row_h_g - 0.12,
             font_size=8.5, color=TEXT_LIGHT)
    add_text(s, state, gap_hx[2] + 0.06, ry + 0.05, gap_hw[2] - 0.12, row_h_g - 0.12,
             font_size=8.5, color=TEXT_GREY)
    add_rect(s, gap_hx[3] + 0.06, ry + 0.18, 1.8, 0.32, pcol)
    add_text(s, prio, gap_hx[3] + 0.06, ry + 0.19, 1.8, 0.28,
             font_size=9, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 11 – Git Workflow & DoD
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Git Workflow & Definition of Done", "Verbindliche Qualitätsgates vor jedem Commit")

# Git-Flow als horizontale Kette
steps = [
    ("git pull", "Immer zuerst\nStand holen", ACCENT_BLUE, "01"),
    ("implement", "Feature /\nBugfix umsetzen", ACCENT_GREEN, "02"),
    ("git add\n<files>", "Nur spezifische\nDateien stagen", ACCENT_TEAL, "03"),
    ("git commit\n-m \"type(scope):\n…\"", "Commit-Format\neinhalten", ACCENT_GOLD, "04"),
    ("git push", "Branch auf\nRemote schieben", ACCENT_RED, "05"),
]
step_y = 1.62
for i, (cmd, sub, col, num) in enumerate(steps):
    sx = 0.35 + i * 2.6
    add_rect(s, sx, step_y, 2.45, 1.65, BG_CARD)
    add_rect(s, sx, step_y, 2.45, 0.055, col)
    add_rect(s, sx + 0.1, step_y + 0.1, 0.35, 0.35, col)
    add_text(s, num, sx + 0.1, step_y + 0.1, 0.35, 0.32,
             font_size=9, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, cmd, sx + 0.55, step_y + 0.1, 1.82, 0.7,
             font_size=10, bold=True, color=TEXT_WHITE)
    add_text(s, sub, sx + 0.1, step_y + 0.9, 2.2, 0.65,
             font_size=8.5, color=TEXT_GREY)
    if i < 4:
        add_text(s, "→", sx + 2.5, step_y + 0.6, 0.22, 0.4,
                 font_size=14, bold=True, color=TEXT_GREY)

# Commit-Format Beispiele
add_rect(s, 0.35, 3.45, 5.9, 2.1, BG_CARD)
add_rect(s, 0.35, 3.45, 0.045, 2.1, ACCENT_GOLD)
add_text(s, "Commit-Format: type(scope): description", 0.5, 3.52, 5.6, 0.38,
         font_size=11, bold=True, color=ACCENT_GOLD)
commit_ex = [
    ("Types:", "feat  fix  refactor  test  docs  chore"),
    ("Scopes:", "Level1–Level6  Arduino  GameManager  UI  NPC  Editor"),
    ("Beispiel:", 'feat(Level3): Helios NPC + Buch-Auswahl implementiert'),
]
for i, (k, v) in enumerate(commit_ex):
    add_text(s, k, 0.5, 4.0 + i * 0.48, 1.0, 0.38,
             font_size=9, bold=True, color=TEXT_GREY)
    add_text(s, v, 1.55, 4.0 + i * 0.48, 4.5, 0.38,
             font_size=9, color=TEXT_LIGHT)

# Definition of Done
add_rect(s, 6.55, 3.45, 6.07, 2.1, BG_CARD)
add_rect(s, 6.55, 3.45, 0.045, 2.1, ACCENT_GREEN)
add_text(s, "Definition of Done", 6.7, 3.52, 5.8, 0.38,
         font_size=11, bold=True, color=ACCENT_GREEN)
dod = [
    "Feature funktioniert im Play Mode (manueller Test)",
    "Unit/Play Tests bestehen",
    "Keine neuen Compiler-Warnungen",
    "Commit mit korrektem Format gepusht",
    "GDD-Gap geschlossen (falls zutreffend)",
]
for i, item in enumerate(dod):
    add_rect(s, 6.7, 4.0 + i * 0.3, 0.22, 0.22, ACCENT_GREEN)
    add_text(s, "✓", 6.71, 4.01 + i * 0.3, 0.2, 0.18,
             font_size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, item, 7.02, 3.99 + i * 0.3, 5.4, 0.28,
             font_size=9, color=TEXT_LIGHT)

# Unten: Warnmeldung
add_rect(s, 0.35, 5.72, 12.7, 1.55, RGBColor(0x1C, 0x15, 0x08))
add_rect(s, 0.35, 5.72, 12.7, 0.045, ACCENT_GOLD)
add_text(s, "⚠  Kein Squash-and-Forget — jede abgeschlossene Aufgabe erhält einen eigenen Commit.",
         0.55, 5.82, 12.2, 0.35, font_size=10, bold=True, color=ACCENT_GOLD)
add_text(s,
         "Niemals `git add -A` oder -A Flag — um versehentliche .env / Binary-Commits zu vermeiden.",
         0.55, 6.22, 12.2, 0.9, font_size=9, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# FOLIE 12 – Nächste Schritte / Sprint-Roadmap
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Nächste Schritte & Sprint-Roadmap", "Priorisierter Entwicklungsplan in drei Phasen")

sprints = [
    ("SPRINT 1", "Level 1 & 2 HIGH-Gaps", ACCENT_RED,
     [
         ("G-01", "Blanket-Hint Objekt\nin 3D-Welt platzieren"),
         ("G-02", "Joshi NPC Narration\nstatt Big Yahu"),
         ("G-03", "Joystick-Kombo ↑↑↓↓\nals Rätsel implementieren"),
     ]),
    ("SPRINT 2", "Level 3 erweitern", ACCENT_GOLD,
     [
         ("G-04", "Helios NPC +\nBuch-Auswahl (Bibel)"),
         ("G-05", "Arduino Color Sensor\nPath integrieren"),
         ("—", "Color Sequence UI\nvia Hardware"),
     ]),
    ("SPRINT 3", "Politur & Release", ACCENT_GREEN,
     [
         ("G-06", "L6 Bunsenbrenner-\nMetapher verbessern"),
         ("QA", "Regression Tests\nfür alle Level"),
         ("🚀", "Build vorbereiten\n& File Hygiene"),
     ]),
]

for i, (sprint, title, col, tasks) in enumerate(sprints):
    x = 0.35 + i * 4.33
    # Sprint-Header
    add_rect(s, x, 1.55, 4.1, 0.85, col)
    add_text(s, sprint, x + 0.15, 1.6, 3.8, 0.38,
             font_size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.15, 1.98, 3.8, 0.32,
             font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Task-Karten
    for j, (tid, desc) in enumerate(tasks):
        ty = 2.55 + j * 1.65
        add_rect(s, x, ty, 4.1, 1.5, BG_CARD)
        add_rect(s, x, ty, 0.04, 1.5, col)
        add_rect(s, x + 0.15, ty + 0.12, 0.55, 0.55, col)
        add_text(s, tid, x + 0.15, ty + 0.12, 0.55, 0.55,
                 font_size=9, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + 0.82, ty + 0.15, 3.1, 1.15,
                 font_size=10, color=TEXT_LIGHT)

# Zeitleiste am Boden
add_rect(s, 0.35, 7.1, 12.7, 0.05, LINE_COLOR)
for i, (label, col) in enumerate([("Sprint 1\n~2 Wochen", ACCENT_RED),
                                   ("Sprint 2\n~2 Wochen", ACCENT_GOLD),
                                   ("Sprint 3\n~1 Woche", ACCENT_GREEN)]):
    x = 0.5 + i * 4.33
    add_rect(s, x, 7.06, 0.06, 0.06, col)
    add_text(s, label, x + 0.12, 7.08, 3.8, 0.32,
             font_size=8, color=col)

# ════════════════════════════════════════════════════════════════════════════
# Speichern
# ════════════════════════════════════════════════════════════════════════════
output = r"c:\Users\flori\Documents\Exit-Game\Exit-Game-Praesentation.pptx"
prs.save(output)
print(f"[OK] Praesentation gespeichert: {output}")
