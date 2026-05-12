"""
Exit-Game — Spielvorstellung & Technische Loesung
Generiert eine 16:9 Praesentation fuer ein Publikum, das das Spiel selbst gespielt hat.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Farbpalette (Dark Theme) ----------
BG_DARK     = RGBColor(0x0D, 0x11, 0x17)
BG_PANEL    = RGBColor(0x16, 0x1B, 0x22)
BG_PANEL_2  = RGBColor(0x1F, 0x26, 0x2F)
BORDER      = RGBColor(0x30, 0x36, 0x3D)
TEXT_MAIN   = RGBColor(0xE6, 0xED, 0xF3)
TEXT_MUTED  = RGBColor(0x8B, 0x94, 0x9E)
ACCENT_GRN  = RGBColor(0x39, 0xD3, 0x53)
ACCENT_BLU  = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_TEA  = RGBColor(0x00, 0xBF, 0xC4)
ACCENT_ORG  = RGBColor(0xFF, 0xA6, 0x57)
ACCENT_RED  = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_VIO  = RGBColor(0xC0, 0x8C, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def set_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_rect(slide, x, y, w, h, fill=BG_PANEL, line=BORDER, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_round_rect(slide, x, y, w, h, fill=BG_PANEL, line=BORDER, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.10
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_MAIN,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_label(slide, x, y, text, color=ACCENT_GRN):
    """small uppercase tag"""
    pad_x = Inches(0.10)
    width = Inches(0.05 + 0.10 * len(text))
    tag = add_rect(slide, x, y, max(width, Inches(0.8)), Inches(0.28), fill=color, line=None)
    add_text(slide, x, y, max(width, Inches(0.8)), Inches(0.28),
             text.upper(), size=9, bold=True, color=BG_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return tag

def slide_header(slide, kicker, title, accent=ACCENT_GRN):
    add_label(slide, Inches(0.6), Inches(0.45), kicker, color=accent)
    add_text(slide, Inches(0.6), Inches(0.80), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=TEXT_MAIN)
    line = add_rect(slide, Inches(0.6), Inches(1.45), Inches(0.6), Pt(3), fill=accent, line=None)
    return line

def add_card(slide, x, y, w, h, *, title=None, body=None,
             accent=ACCENT_BLU, fill=BG_PANEL):
    add_rect(slide, x, y, w, h, fill=fill, line=BORDER)
    # accent strip
    add_rect(slide, x, y, Inches(0.07), h, fill=accent, line=None)
    inner_x = x + Inches(0.20)
    inner_w = w - Inches(0.30)
    if title:
        add_text(slide, inner_x, y + Inches(0.12), inner_w, Inches(0.35),
                 title, size=13, bold=True, color=TEXT_MAIN)
    if body:
        add_text(slide, inner_x, y + Inches(0.50), inner_w, h - Inches(0.55),
                 body, size=11, color=TEXT_MUTED)

def add_chip(slide, x, y, text, color=ACCENT_BLU, w=None):
    width = w or Inches(0.10 + 0.085 * len(text))
    add_round_rect(slide, x, y, width, Inches(0.32), fill=BG_PANEL_2, line=color, line_w=0.75)
    add_text(slide, x, y, width, Inches(0.32), text, size=10, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return width

def footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "Exit-Game  -  Spielvorstellung & Technische Loesung",
             size=10, color=TEXT_MUTED)
    add_text(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

TOTAL = 14

# ============================================================
# SLIDE 1 - TITEL
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)

# Deko: diagonale Akzentstreifen
add_rect(s, Inches(-1.0), Inches(6.5), Inches(20), Inches(0.10), fill=ACCENT_GRN, line=None)
add_rect(s, Inches(-1.0), Inches(6.7), Inches(20), Inches(0.04), fill=ACCENT_BLU, line=None)

# Eyebrow
add_label(s, Inches(0.9), Inches(1.4), "EXIT-GAME PROJEKT", color=ACCENT_GRN)
add_text(s, Inches(0.85), Inches(1.85), Inches(12), Inches(1.4),
         "Vom realen Raum\nins digitale Spiel",
         size=58, bold=True, color=TEXT_MAIN)
add_text(s, Inches(0.9), Inches(4.05), Inches(12), Inches(0.6),
         "Sechs Raetsel - sechs Mechaniken - eine Mission: Entkommen aus dem Knast.",
         size=18, color=TEXT_MUTED)

# Author cards
add_round_rect(s, Inches(0.9), Inches(5.2), Inches(3.6), Inches(1.05), fill=BG_PANEL, line=BORDER)
add_text(s, Inches(1.1), Inches(5.35), Inches(3.2), Inches(0.4),
         "Florian Maier", size=14, bold=True, color=TEXT_MAIN)
add_text(s, Inches(1.1), Inches(5.7), Inches(3.2), Inches(0.4),
         "Entwicklung & Gamedesign", size=11, color=TEXT_MUTED)

add_round_rect(s, Inches(4.7), Inches(5.2), Inches(3.6), Inches(1.05), fill=BG_PANEL, line=BORDER)
add_text(s, Inches(4.9), Inches(5.35), Inches(3.2), Inches(0.4),
         "Unity + Arduino", size=14, bold=True, color=TEXT_MAIN)
add_text(s, Inches(4.9), Inches(5.7), Inches(3.2), Inches(0.4),
         "Hybrider Hard- & Software-Stack", size=11, color=TEXT_MUTED)

add_round_rect(s, Inches(8.5), Inches(5.2), Inches(3.6), Inches(1.05), fill=BG_PANEL, line=BORDER)
add_text(s, Inches(8.7), Inches(5.35), Inches(3.2), Inches(0.4),
         "Mai 2026", size=14, bold=True, color=TEXT_MAIN)
add_text(s, Inches(8.7), Inches(5.7), Inches(3.2), Inches(0.4),
         "Projektabschluss & Demo", size=11, color=TEXT_MUTED)

footer(s, 1, TOTAL)

# ============================================================
# SLIDE 2 - HERZLICH WILLKOMMEN / DANKE FUERS SPIELEN
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Willkommen", "Sie haben das Spiel gerade gespielt - was nun?", accent=ACCENT_GRN)

add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.7),
         "Diese Praesentation gibt Ihnen einen Blick hinter die Kulissen.",
         size=18, color=TEXT_MUTED)

# Drei Spalten: was Sie erlebt haben / was Sie sehen werden / was Sie mitnehmen
col_y = Inches(2.7); col_h = Inches(3.6); col_w = Inches(4.0)
cols = [
    ("Was Sie erlebt haben",
     "Sechs Raetsel-Stationen, gesteuert ueber Tastatur, Maus und echte Hardware-Sensoren.",
     ACCENT_GRN, "01"),
    ("Was Sie jetzt sehen",
     "Den Ablauf hinter dem Vorhang: Architektur, Arduino-Bruecke, Level-Mechaniken im Detail.",
     ACCENT_BLU, "02"),
    ("Was Sie mitnehmen",
     "Wie reale Hardware und Unity Hand in Hand spielen - und wo die Knackpunkte lagen.",
     ACCENT_TEA, "03"),
]
for i, (t, b, c, n) in enumerate(cols):
    x = Inches(0.6 + i * 4.2)
    add_rect(s, x, col_y, col_w, col_h, fill=BG_PANEL, line=BORDER)
    add_rect(s, x, col_y, col_w, Inches(0.10), fill=c, line=None)
    add_text(s, x + Inches(0.3), col_y + Inches(0.3), col_w - Inches(0.6), Inches(0.5),
             n, size=32, bold=True, color=c)
    add_text(s, x + Inches(0.3), col_y + Inches(1.05), col_w - Inches(0.6), Inches(0.5),
             t, size=16, bold=True, color=TEXT_MAIN)
    add_text(s, x + Inches(0.3), col_y + Inches(1.55), col_w - Inches(0.6), Inches(2.0),
             b, size=12, color=TEXT_MUTED)

footer(s, 2, TOTAL)

# ============================================================
# SLIDE 3 - SPIELKONZEPT & STORY
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Konzept", "Das Spiel in einem Satz", accent=ACCENT_GRN)

# Hero quote
add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.5), fill=BG_PANEL, line=BORDER)
add_rect(s, Inches(0.6), Inches(1.7), Inches(0.10), Inches(1.5), fill=ACCENT_GRN, line=None)
add_text(s, Inches(1.0), Inches(1.9), Inches(11.5), Inches(1.2),
         '"Ein Escape Room im Digitalen - mit echten Knoepfen,\n'
         ' Sensoren und einem digitalen Mithaeftling."',
         size=20, bold=True, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)

# Story / Genre / Plattform / Spieldauer
cards = [
    ("Setting", "Gefaengnis", "Zelle, Lager, Bibliothek, Computerraum, Werkstatt, Tor.", ACCENT_GRN),
    ("Genre", "Escape / Puzzle", "Erstperson-3D & UI-Panels, sechs Raetsel-Mechaniken.", ACCENT_BLU),
    ("Plattform", "PC + Hardware", "Unity Windows-Build, Arduino-Bridge ueber USB.", ACCENT_TEA),
    ("Spieldauer", "~15-25 Min", "Lineare Progression, eine Sitzung, kein Speicherstand.", ACCENT_ORG),
]
for i, (label, title, body, color) in enumerate(cards):
    x = Inches(0.6 + i * 3.05); y = Inches(3.5); w = Inches(2.95); h = Inches(2.9)
    add_rect(s, x, y, w, h, fill=BG_PANEL, line=BORDER)
    add_label(s, x + Inches(0.2), y + Inches(0.25), label, color=color)
    add_text(s, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(0.6),
             title, size=20, bold=True, color=TEXT_MAIN)
    add_text(s, x + Inches(0.2), y + Inches(1.4), w - Inches(0.4), Inches(1.5),
             body, size=11, color=TEXT_MUTED)

footer(s, 3, TOTAL)

# ============================================================
# SLIDE 4 - SPIELABLAUF (VORGANG): DER ROTE FADEN
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Vorgang", "So spielt sich das Spiel - von Start bis Tor", accent=ACCENT_BLU)

# Horizontaler Fluss mit 6 Stationen
y_track = Inches(2.5)
# Verbindungslinie
add_rect(s, Inches(1.05), y_track + Inches(0.95), Inches(11.2), Pt(2), fill=ACCENT_BLU, line=None)

stations = [
    ("L1", "Zelle", "Hex-Code\nauf Decke", ACCENT_GRN),
    ("L2", "Lager", "Staubwand\n+ Sensor", ACCENT_BLU),
    ("L3", "Biblio", "Bibel +\nFarbcode", ACCENT_TEA),
    ("L4", "Computer", "Stealth-\nMinigame", ACCENT_ORG),
    ("L5", "Werkstatt", "Breadboard-\nSchaltung", ACCENT_VIO),
    ("L6", "Tor", "Hitze /\nTemperatur", ACCENT_RED),
]
for i, (lvl, title, body, c) in enumerate(stations):
    cx = Inches(1.4 + i * 1.85)
    # Knoten
    add_round_rect(s, cx, y_track + Inches(0.55), Inches(0.8), Inches(0.8), fill=BG_PANEL, line=c, line_w=2.0)
    add_text(s, cx, y_track + Inches(0.55), Inches(0.8), Inches(0.8),
             lvl, size=18, bold=True, color=c, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Label oben
    add_text(s, cx - Inches(0.4), y_track - Inches(0.05), Inches(1.6), Inches(0.4),
             title, size=14, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
    # Body unten
    add_text(s, cx - Inches(0.5), y_track + Inches(1.55), Inches(1.8), Inches(0.9),
             body, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Unten: Mechanik-Erklaerung
add_rect(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.5), fill=BG_PANEL, line=BORDER)
add_rect(s, Inches(0.6), Inches(5.2), Inches(0.10), Inches(1.5), fill=ACCENT_BLU, line=None)
add_text(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.4),
         "Linearer Levelfluss mit klarer Eskalation der Mechanik",
         size=14, bold=True, color=TEXT_MAIN)
add_text(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.9),
         "Jedes Level steigert die Komplexitaet: einfache Code-Eingabe -> Sensorinteraktion -> Stealth -> "
         "physische Schaltung -> finale Sensor-Verifikation. Big Yahu (NPC) fuehrt narrativ durch die Stationen.",
         size=11, color=TEXT_MUTED)

footer(s, 4, TOTAL)

# ============================================================
# SLIDE 5 - DIE SECHS LEVEL IM DETAIL (TABELLE)
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Level im Detail", "Was passiert wo - was war die Herausforderung?", accent=ACCENT_BLU)

# Tabellenkopf
table_x = Inches(0.6); table_y = Inches(1.85)
col_w = [Inches(0.7), Inches(1.9), Inches(3.8), Inches(3.3), Inches(2.4)]
total_w = sum(col_w, Inches(0))

# Header row
hx = table_x
add_rect(s, table_x, table_y, total_w, Inches(0.45), fill=BG_PANEL_2, line=BORDER)
headers = ["#", "Schauplatz", "Mechanik", "Spielereingabe", "Hardware"]
for i, h in enumerate(headers):
    add_text(s, hx + Inches(0.1), table_y + Inches(0.05), col_w[i] - Inches(0.2), Inches(0.35),
             h, size=11, bold=True, color=ACCENT_BLU)
    hx += col_w[i]

# Rows
rows = [
    ("1", "Zelle", "Hex-Code 66A auf Decke -> Numpad", "Tastatur / Numpad-UI", "-"),
    ("2", "Abstellraum", "Staub wegpusten + Joystick-Combo", "Mikrofon-Pegel / Tasten", "Feuchtesensor (geplant)"),
    ("3", "Bibliothek", "Richtige Bibel + Farbcode scannen", "Buchauswahl + Farbsensor", "Color-Sensor"),
    ("4", "Computer", "Stealth-Minigame, Wachen umgehen", "Maus / WASD im UI-Raster", "-"),
    ("5", "Werkstatt", "Breadboard-Schaltung loesen", "Echte Kabel & Komponenten", "Arduino I/O Pins"),
    ("6", "Tor", "Hitze ans Schloss - Foehn/Brenner", "Temperatur halten", "Temperatursensor"),
]
for r_i, row in enumerate(rows):
    ry = table_y + Inches(0.45) + Inches(r_i * 0.6)
    bg = BG_PANEL if r_i % 2 == 0 else BG_PANEL_2
    add_rect(s, table_x, ry, total_w, Inches(0.6), fill=bg, line=BORDER, line_w=0.4)
    rx = table_x
    for ci, val in enumerate(row):
        if ci == 0:
            add_text(s, rx + Inches(0.1), ry + Inches(0.13), col_w[ci] - Inches(0.2), Inches(0.4),
                     val, size=14, bold=True, color=ACCENT_GRN)
        else:
            add_text(s, rx + Inches(0.1), ry + Inches(0.15), col_w[ci] - Inches(0.2), Inches(0.4),
                     val, size=11, color=TEXT_MAIN if ci == 1 else TEXT_MUTED)
        rx += col_w[ci]

footer(s, 5, TOTAL)

# ============================================================
# SLIDE 6 - TECHNISCHE LOESUNG: GROSSARCHITEKTUR
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Technische Loesung", "Die drei Schichten der Architektur", accent=ACCENT_TEA)

# Drei horizontale Schichten
def layer(y, title, sub, color, items):
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.5), fill=BG_PANEL, line=BORDER)
    add_rect(s, Inches(0.6), y, Inches(0.10), Inches(1.5), fill=color, line=None)
    add_label(s, Inches(0.9), y + Inches(0.15), title, color=color)
    add_text(s, Inches(0.9), y + Inches(0.55), Inches(4.5), Inches(0.4),
             sub, size=14, bold=True, color=TEXT_MAIN)
    # Chips rechts
    cx = Inches(5.5); cy = y + Inches(0.55)
    for chip, ccolor in items:
        w = add_chip(s, cx, cy, chip, color=ccolor)
        cx = cx + w + Inches(0.12)
        if cx > Inches(12.4):
            cx = Inches(5.5); cy = cy + Inches(0.45)

layer(Inches(1.85), "Praesentation",
      "Unity 3D, UI-Panels, NPCs",
      ACCENT_GRN,
      [("Unity 2022", ACCENT_GRN), ("Input System", ACCENT_GRN), ("TextMeshPro", ACCENT_GRN),
       ("Cinemachine", ACCENT_GRN), ("BigYahuDialogSystem", ACCENT_BLU)])

layer(Inches(3.55), "Logik / Steuerung",
      "C# MonoBehaviours, GameManager",
      ACCENT_BLU,
      [("GameManager", ACCENT_BLU), ("Level1_Cell", ACCENT_BLU), ("Level2_DustWall", ACCENT_BLU),
       ("Level3_ColorCode", ACCENT_BLU), ("Level4_Stealth", ACCENT_BLU),
       ("Level5_Breadboard", ACCENT_BLU), ("Level6_FinalGate", ACCENT_BLU)])

layer(Inches(5.25), "Hardware-Bruecke",
      "Arduino, Sensoren, Serial-Protokoll",
      ACCENT_TEA,
      [("ArduinoBridge", ACCENT_TEA), ("SerialPort", ACCENT_TEA), ("TCP-Emulator", ACCENT_TEA),
       ("Color-Sensor", ACCENT_ORG), ("Humidity", ACCENT_ORG),
       ("Temperature", ACCENT_ORG), ("Breadboard I/O", ACCENT_ORG)])

footer(s, 6, TOTAL)

# ============================================================
# SLIDE 7 - DATENFLUSS DURCH DAS SYSTEM
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Datenfluss", "Von der Hand am Sensor bis zum Levelwechsel", accent=ACCENT_TEA)

# 5 Boxen mit Pfeilen
boxes = [
    ("Sensor", "Spieler\ninteragiert", ACCENT_ORG, "Drehknopf, Foehn,\nKabel, Mikro"),
    ("Arduino", "Misst &\nsendet", ACCENT_TEA, "0x60:42\\n\nueber Serial"),
    ("Bridge", "Liest auf\nBG-Thread", ACCENT_BLU, "ConcurrentQueue\nmainthread-safe"),
    ("Level-Skript", "Parst &\nbewertet", ACCENT_GRN, "OnTemp / OnColor\n/ OnHumidity"),
    ("GameManager", "Schaltet\nLevel um", ACCENT_VIO, "CompleteCurrentLevel\n-> Panel-Wechsel"),
]
box_y = Inches(2.2); box_h = Inches(2.6); box_w = Inches(2.05); gap = Inches(0.35)
for i, (label, title, color, body) in enumerate(boxes):
    x = Inches(0.6) + i * (box_w + gap)
    add_rect(s, x, box_y, box_w, box_h, fill=BG_PANEL, line=BORDER)
    add_rect(s, x, box_y, box_w, Inches(0.10), fill=color, line=None)
    add_label(s, x + Inches(0.2), box_y + Inches(0.25), label, color=color)
    add_text(s, x + Inches(0.2), box_y + Inches(0.7), box_w - Inches(0.4), Inches(0.9),
             title, size=15, bold=True, color=TEXT_MAIN)
    add_text(s, x + Inches(0.2), box_y + Inches(1.55), box_w - Inches(0.4), Inches(1.0),
             body, size=10, color=TEXT_MUTED, font="Consolas")
    # Pfeil zur naechsten Box
    if i < len(boxes) - 1:
        arrow_x = x + box_w + Inches(0.02)
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, box_y + Inches(1.1),
                               Inches(0.3), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = ACCENT_TEA
        a.line.fill.background()
        a.shadow.inherit = False

# Unten: Highlight
add_rect(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.4), fill=BG_PANEL, line=BORDER)
add_rect(s, Inches(0.6), Inches(5.3), Inches(0.10), Inches(1.4), fill=ACCENT_TEA, line=None)
add_text(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(0.4),
         "Thread-sicher dank ConcurrentQueue + Mainthread-Dispatch",
         size=14, bold=True, color=TEXT_MAIN)
add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.8),
         "Serial-Reads laufen auf einem dedizierten Hintergrund-Thread - alle Events werden ueber eine "
         "ConcurrentQueue auf den Unity-Mainthread geleitet. Kein Freeze, kein Race-Condition, sauberer "
         "Reconnect mit Timeout.",
         size=11, color=TEXT_MUTED)

footer(s, 7, TOTAL)

# ============================================================
# SLIDE 8 - ARDUINO-BRIDGE IM DETAIL
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Arduino-Bridge", "Das Herz der Hybrid-Loesung", accent=ACCENT_ORG)

# Links: Code-Snippet
add_rect(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.85), fill=BG_PANEL_2, line=BORDER)
add_label(s, Inches(0.8), Inches(1.9), "Protokoll", color=ACCENT_ORG)

code = (
    "// Eingehende Nachrichten (Arduino -> Unity)\n"
    "KEY:5         -> Keypad-Taste 5\n"
    "HUMIDITY:72   -> Feuchtigkeit 72%\n"
    "COLOR:#FF8800 -> Farbsensor Hex\n"
    "TEMP:58       -> Temperatur in Grad\n"
    "60:58         -> Hex-Format (Cmd 0x60)\n"
    "\n"
    "// Ausgehende Nachrichten (Unity -> Arduino)\n"
    "Send(0xFF, \"ping\");      // keepalive\n"
    "Send(0x50, \"reset\");     // breadboard\n"
    "\n"
    "// Doppelte API in C#\n"
    "OnKeypadKey += key => ...;\n"
    "OnHumidity  += val => ...;\n"
    "RegisterHandler(0x60, payload => ...);"
)
add_text(s, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.2),
         code, size=11, color=TEXT_MAIN, font="Consolas")

# Rechts: Features
add_label(s, Inches(7.0), Inches(1.85), "Features", color=ACCENT_GRN)
feats = [
    ("Auto-Reconnect",   "Timeout + reconnectDelay - Hardware kann ein- und ausgesteckt werden", ACCENT_GRN),
    ("TCP-Emulator",     "Spiel laeuft komplett ohne echtes Arduino - Python-Emulator simuliert Sensoren", ACCENT_BLU),
    ("Fallback-Eingabe", "Jedes Hardware-Raetsel hat einen Tastatur/Maus-Plan-B im Inspector aktivierbar", ACCENT_TEA),
    ("Singleton + DontDestroyOnLoad", "Verbindung bleibt scenenuebergreifend offen", ACCENT_ORG),
]
y = Inches(2.3)
for title, body, c in feats:
    add_rect(s, Inches(7.0), y, Inches(5.7), Inches(1.0), fill=BG_PANEL, line=BORDER)
    add_rect(s, Inches(7.0), y, Inches(0.08), Inches(1.0), fill=c, line=None)
    add_text(s, Inches(7.2), y + Inches(0.10), Inches(5.4), Inches(0.4),
             title, size=13, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(7.2), y + Inches(0.48), Inches(5.4), Inches(0.5),
             body, size=11, color=TEXT_MUTED)
    y += Inches(1.10)

footer(s, 8, TOTAL)

# ============================================================
# SLIDE 9 - LEVEL-TECHNIK 1: ZELLE & STAUB
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Level 1 & 2", "Zelle und Abstellraum - der Einstieg", accent=ACCENT_GRN)

# Zwei grosse Karten
def level_card(x, y, w, h, lvl, title, mechanik, scripts, hardware, color):
    add_rect(s, x, y, w, h, fill=BG_PANEL, line=BORDER)
    add_rect(s, x, y, w, Inches(0.10), fill=color, line=None)
    add_label(s, x + Inches(0.3), y + Inches(0.25), lvl, color=color)
    add_text(s, x + Inches(0.3), y + Inches(0.7), w - Inches(0.6), Inches(0.5),
             title, size=20, bold=True, color=TEXT_MAIN)

    yy = y + Inches(1.4)
    add_text(s, x + Inches(0.3), yy, Inches(1.4), Inches(0.4),
             "Mechanik", size=10, bold=True, color=color)
    add_text(s, x + Inches(1.6), yy, w - Inches(1.9), Inches(0.4),
             mechanik, size=11, color=TEXT_MAIN)
    yy += Inches(0.55)
    add_text(s, x + Inches(0.3), yy, Inches(1.4), Inches(0.4),
             "Skripte", size=10, bold=True, color=color)
    add_text(s, x + Inches(1.6), yy, w - Inches(1.9), Inches(0.8),
             scripts, size=11, color=TEXT_MUTED, font="Consolas")
    yy += Inches(0.95)
    add_text(s, x + Inches(0.3), yy, Inches(1.4), Inches(0.4),
             "Hardware", size=10, bold=True, color=color)
    add_text(s, x + Inches(1.6), yy, w - Inches(1.9), Inches(0.4),
             hardware, size=11, color=TEXT_MAIN)

level_card(Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 1", "Die Zelle",
           "Hex-Code 66A (->1642) auf Decke. Numpad-Eingabe oeffnet das Schloss.",
           "Level1_Cell.cs\nToiletInteraction.cs\nNumpadController.cs",
           "Keine - reines Unity-Raetsel mit 3D-Welt",
           ACCENT_GRN)

level_card(Inches(6.75), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 2", "Der Abstellraum",
           "Staub von der Wand pusten + Joystick-Kombination Pfeil hoch-hoch-runter-runter.",
           "Level2_DustWall.cs",
           "Feuchtigkeitssensor (Atemerkennung) - geplant, mit Tastatur-Fallback",
           ACCENT_BLU)

footer(s, 9, TOTAL)

# ============================================================
# SLIDE 10 - LEVEL-TECHNIK 2: BIBLIO & COMPUTER
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Level 3 & 4", "Bibliothek und Computerraum - die Eskalation", accent=ACCENT_TEA)

level_card(Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 3", "Die Bibliothek",
           "Richtige Bibel finden, dann Farbsequenz mit Sensor abscannen.",
           "Level3_ColorCode.cs\nHeliosInteraction.cs",
           "Farbsensor (TCS34725) - sendet COLOR:#hex an Unity",
           ACCENT_TEA)

level_card(Inches(6.75), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 4", "Der Computerraum",
           "Stealth-Minigame: Spieler-Token muss Wachen mit Sichtkegeln umgehen.",
           "Level4_StealthMinigame.cs\nBuildLevel4Computer.cs (Editor)",
           "Keine - reines UI-Minigame mit Input System",
           ACCENT_ORG)

footer(s, 10, TOTAL)

# ============================================================
# SLIDE 11 - LEVEL-TECHNIK 3: WERKSTATT & TOR
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Level 5 & 6", "Werkstatt und Tor - das Finale", accent=ACCENT_VIO)

level_card(Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 5", "Die Werkstatt",
           "Reale Breadboard-Schaltung mit Widerstaenden, LED und Schalter. Korrekte Verschaltung schliesst den Stromkreis.",
           "Level5_Breadboard.cs\nBuildLevel5Breadboard.cs",
           "Arduino Digital-/Analog-Pins (Cmd 0x50) - echtes Steckbrett",
           ACCENT_VIO)

level_card(Inches(6.75), Inches(1.75), Inches(6.0), Inches(4.7),
           "Level 6", "Das Tor",
           "Schloss erhitzen (Foehn/Bunsenbrenner-Metapher) - Temperatur muss in Schwelle gehalten werden.",
           "Level6_FinalGate.cs\nBuildLevel6FinalGate.cs",
           "Temperatursensor (DS18B20) - sendet TEMP:xx an Unity (Cmd 0x60)",
           ACCENT_RED)

footer(s, 11, TOTAL)

# ============================================================
# SLIDE 12 - HERAUSFORDERUNGEN & LOESUNGEN
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Lessons Learned", "Was war hart - und wie haben wir es geloest?", accent=ACCENT_ORG)

challenges = [
    ("Serial blockiert Unity", ACCENT_RED,
     "SerialPort.Read() ist synchron und friert den Mainthread ein.",
     "Hintergrund-Thread + ConcurrentQueue + Mainthread-Dispatch im Update().",
     ACCENT_GRN),
    ("Hardware nicht immer da", ACCENT_RED,
     "Tests, Demos und Praesentationen ohne echtes Arduino moeglich machen.",
     "TCP-Emulator (Python) + Inspector-Schalter useTcpEmulator + Fallback-Mode pro Level.",
     ACCENT_GRN),
    ("Szenenwechsel reisst Singletons", ACCENT_RED,
     "ArduinoBridge und GameManager duerfen nicht neu initialisiert werden.",
     "DontDestroyOnLoad + Instance-Check in Awake() - alte Instanzen werden zerstoert.",
     ACCENT_GRN),
    ("Editor-Builder serialisieren nicht", ACCENT_RED,
     "Reflection-basiertes Setzen privater Felder ueberlebte das Speichern nicht.",
     "Umstellung auf SerializedObject + FindProperty + ApplyModifiedPropertiesWithoutUndo.",
     ACCENT_GRN),
]
y = Inches(1.85); row_h = Inches(1.20)
for i, (title, c1, problem, solution, c2) in enumerate(challenges):
    add_rect(s, Inches(0.6), y, Inches(12.1), row_h, fill=BG_PANEL, line=BORDER)
    # Title
    add_label(s, Inches(0.8), y + Inches(0.15), "Problem", color=c1)
    add_text(s, Inches(0.8), y + Inches(0.5), Inches(4.5), Inches(0.6),
             title, size=14, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(0.8), y + Inches(0.85), Inches(4.5), Inches(0.4),
             problem, size=10, color=TEXT_MUTED)
    # Arrow
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.4), y + Inches(0.45),
                           Inches(0.4), Inches(0.35))
    a.fill.solid(); a.fill.fore_color.rgb = ACCENT_TEA
    a.line.fill.background()
    a.shadow.inherit = False
    # Solution
    add_label(s, Inches(5.95), y + Inches(0.15), "Loesung", color=c2)
    add_text(s, Inches(5.95), y + Inches(0.5), Inches(6.6), Inches(0.7),
             solution, size=12, color=TEXT_MAIN)
    y += row_h + Inches(0.08)

footer(s, 12, TOTAL)

# ============================================================
# SLIDE 13 - AUSBLICK / WAS KOMMT ALS NAECHSTES
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s, "Ausblick", "Naechste Schritte fuer das Projekt", accent=ACCENT_GRN)

# Roadmap-Spalten
cols = [
    ("KURZFRISTIG", ACCENT_GRN, "naechste 2 Wochen", [
        "Joshi-NPC mit eigener Stimme in Level 2",
        "Bibel-Auswahl-Mechanik in Level 3",
        "Decken-Hex-Hinweis 66A als 3D-Objekt",
        "Fallback-UI fuer alle Sensoren testen",
    ]),
    ("MITTELFRISTIG", ACCENT_BLU, "naechster Monat", [
        "Soundkulisse pro Level (Ambient + SFX)",
        "Speicher-/Resume-System fuer lange Sessions",
        "Verbessertes Stealth-AI in Level 4",
        "Foehn/Brenner-Visualisierung in Level 6",
    ]),
    ("LANGFRISTIG", ACCENT_TEA, "Ausblick", [
        "Mehrspieler-Variante (zwei Spieler kooperieren)",
        "Editor zum eigenen Level-Erstellen",
        "Web-Build mit emulierter Hardware",
        "Showcase auf Event mit echtem Aufbau",
    ]),
]
for i, (label, color, sub, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.2); y = Inches(1.85); w = Inches(4.0); h = Inches(4.8)
    add_rect(s, x, y, w, h, fill=BG_PANEL, line=BORDER)
    add_rect(s, x, y, w, Inches(0.10), fill=color, line=None)
    add_label(s, x + Inches(0.3), y + Inches(0.25), label, color=color)
    add_text(s, x + Inches(0.3), y + Inches(0.65), w - Inches(0.6), Inches(0.4),
             sub, size=11, color=TEXT_MUTED)
    yy = y + Inches(1.15)
    for it in items:
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), yy + Inches(0.12),
                                    Inches(0.12), Inches(0.12))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()
        bullet.shadow.inherit = False
        add_text(s, x + Inches(0.55), yy, w - Inches(0.85), Inches(0.7),
                 it, size=11, color=TEXT_MAIN)
        yy += Inches(0.75)

footer(s, 13, TOTAL)

# ============================================================
# SLIDE 14 - DANKE / Q&A
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s)

# Vollflaechige Grafik
add_rect(s, Inches(-1.0), Inches(0.5), Inches(20), Inches(0.04), fill=ACCENT_GRN, line=None)
add_rect(s, Inches(-1.0), Inches(7.0), Inches(20), Inches(0.04), fill=ACCENT_BLU, line=None)

add_label(s, Inches(0.9), Inches(1.5), "DANKE", color=ACCENT_GRN)
add_text(s, Inches(0.85), Inches(2.0), Inches(12), Inches(1.4),
         "Vielen Dank fuers\nMitspielen und Zuhoeren.",
         size=58, bold=True, color=TEXT_MAIN)

add_text(s, Inches(0.9), Inches(4.4), Inches(12), Inches(0.6),
         "Fragen, Feedback, Ausprobieren?",
         size=22, color=ACCENT_BLU, bold=True)

# Q&A Hinweise
qa = [
    ("Spielerlebnis", "Was war intuitiv, was war frustrierend?"),
    ("Technik",       "Wie soll sich das Hardware-Spiel anfuehlen?"),
    ("Ideen",         "Welche Mechanik fehlt - was wuerden Sie hinzufuegen?"),
]
for i, (t, b) in enumerate(qa):
    x = Inches(0.9 + i * 4.05); y = Inches(5.2)
    add_round_rect(s, x, y, Inches(3.8), Inches(1.4), fill=BG_PANEL, line=BORDER)
    add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.4), Inches(0.4),
             t, size=14, bold=True, color=ACCENT_GRN)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(3.4), Inches(0.7),
             b, size=11, color=TEXT_MUTED)

footer(s, 14, TOTAL)

# ============================================================
# SAVE
# ============================================================
out = r"c:\Users\flori\Documents\Exit-Game\Exit-Game-Spielvorstellung.pptx"
prs.save(out)
print(f"[OK] Praesentation gespeichert: {out}")
print(f"     {TOTAL} Folien, 16:9 Dark Theme")
